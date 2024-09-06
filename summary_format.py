from fpdf import FPDF
from docx import Document
from pptx import Presentation
import tempfile

class OutputFormat():
    def pdf(self, summary, output_io):
        with tempfile.NamedTemporaryFile(delete=False, suffix='.pdf') as temp_pdf:
            pdf = FPDF()
            pdf.add_page()
            pdf.set_font("Arial", size=12)
            pdf.multi_cell(0, 10, summary)
            pdf.output(temp_pdf.name)
            temp_pdf.seek(0)
            output_io.write(temp_pdf.read())
    
    def word(self, summary, output_io):
        doc = Document()
        doc.add_paragraph(summary)
        doc.save(output_io)
    
    def text(self, summary, output_io):
        output_io.write(summary.encode('utf-8'))
    
    def split_text(self, summary, max_length):
        words = summary.split()
        chunks = []
        current_chunk = ""
        for word in words:
            if len(current_chunk) + len(word) + 1 > max_length:
                chunks.append(current_chunk)
                current_chunk = word
            else:
                if current_chunk:
                    current_chunk += " "
                current_chunk += word

        if current_chunk:
            chunks.append(current_chunk)
    
        return chunks
    
    def powerpoint(self, summary, output_io):
        prs = Presentation()
        chunks = self.split_text(summary, 20)
        for i, chunk in enumerate(chunks):
            slide_layout = prs.slide_layouts[1]  
            slide = prs.slides.add_slide(slide_layout)
            title = slide.shapes.title
            title.text = f"Slide {i + 1}"
            content = slide.placeholders[1]
            content.text = chunk
        prs.save(output_io)
