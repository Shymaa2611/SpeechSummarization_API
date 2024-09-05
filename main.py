from fastapi import FastAPI
from fpdf import FPDF
from docx import Document
from pptx import Presentation
from pydub import AudioSegment
from transformers import pipeline
from transformers import GPT2LMHeadModel, GPT2Tokenizer
from summarize_text import summarize_text
from fastapi import FastAPI, File, UploadFile, Form
from fastapi.responses import StreamingResponse
from io import BytesIO
from fastapi.templating import Jinja2Templates
from fastapi import Request

app = FastAPI()
templates = Jinja2Templates(directory="templates")

class Speech2Text():
    def audio_segments(self,audio_url, segment_length_ms=15000):
      audio = AudioSegment.from_mp3(audio_url)
      total_length_ms = len(audio)
      segments = []
      for i in range(0, total_length_ms, segment_length_ms):
        segment = audio[i:i+segment_length_ms]
        segments.append(segment)
        print(f'Segment {i // segment_length_ms + 1} processed')
      return segments

    def load_model(self):
          pipe = pipeline("automatic-speech-recognition", model="openai/whisper-large-v3")
          return pipe
    
    def speech2text(self,audio_url):
       pipe=self.load_model()
       segments=self.audio_segments(audio_url)
       full_text=""
       for segment in segments:
          text=pipe(segment)
          full_text+=text
       return full_text
       
class TextSummarization():
    def load_model(self):
       model = GPT2LMHeadModel.from_pretrained('./checkpoint')
       tokenizer = GPT2Tokenizer.from_pretrained('./checkpoint')
       return model,tokenizer
    
       
    def text2summarize(self,full_text):
        model,tokenizer=self.load_model()
        summary=summarize_text(full_text,tokenizer,model)
        return summary

class Speech2TextSummarization():
   def __init__(self):
       self.s2t=Speech2Text()
       self.ts=TextSummarization()

   def speech2textsummarization(self,audio_url):
       text=self.s2t(audio_url)
       summarize=self.ts(text)
       return summarize
   
class OutputFormat():
    def pdf(self,summary):
       pdf = FPDF()
       pdf.add_page()
       pdf.set_font("Arial", size=12)
       pdf.multi_cell(0, 10, summary)
       pdf_output = "output.pdf"  
       pdf.output(pdf_output)
       print(f"PDF generated: {pdf_output}")

    def word(self,summary):
      doc = Document()
      doc.add_paragraph(summary)
      output_file="output.doc"
      doc.save(output_file)
      print(f"Word document generated: {output_file}")

    def text(self,summary):
      file_path="output.txt"
      with open(file_path, 'w') as file:
        file.write(summary)
      print(f"Text has been written to {file_path}")
    
    def split_text(self,summary, max_length):
      words = summary.split()
      chunks = []
      current_chunk = ""
      for word in words:
        if len(current_chunk) + len(word) + 1 > max_length:
            chunks.append(current_chunk)
            current_chunk = word
        else:
            if current_chunk:
                current_chunk += " "
            current_chunk += word

      if current_chunk:
        chunks.append(current_chunk)
    
      return chunks
       
       
    def powerpoint(self,summary):
      prs = Presentation()
      file_path="output.pptx"
      chunks =self.split_text(summary,20)
      for i, chunk in enumerate(chunks):
        slide_layout = prs.slide_layouts[1]  
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = f"Slide {i + 1}"
        content = slide.placeholders[1]
        content.text = chunk
    
      prs.save(file_path)
      print(f"PowerPoint presentation saved to {file_path}")

       
@app.get("/")
async def main_page(request: Request):
    return templates.TemplateResponse("index.html", {"request": request})

@app.post("/upload/")
async def upload_audio(file: UploadFile = File(...), format: str = Form(...)):
    file_bytes = await file.read()
    audio_file = BytesIO(file_bytes)
    audio_path = 'temp_audio.mp3'
    with open(audio_path, 'wb') as f:
        f.write(file_bytes)
    s2ts = Speech2TextSummarization()
    summary = s2ts.speech2textsummarization(audio_path)
    output_format = OutputFormat()
    if format == 'pdf':
            pdf_io = BytesIO()
            output_format.pdf(summary)
            pdf_io.seek(0)
            return StreamingResponse(pdf_io, media_type='application/pdf', headers={"Content-Disposition": "attachment; filename=summary.pdf"})
    elif format == 'word':
            doc_io = BytesIO()
            output_format.word(summary)
            doc_io.seek(0)
            return StreamingResponse(doc_io, media_type='application/vnd.openxmlformats-officedocument.wordprocessingml.document', headers={"Content-Disposition": "attachment; filename=summary.docx"})     
    elif format == 'text':
            text_io = BytesIO()
            output_format.text(summary)
            text_io.seek(0)
            return StreamingResponse(text_io, media_type='text/plain', headers={"Content-Disposition": "attachment; filename=summary.txt"})
        
    elif format == 'pptx':
            ppt_io = BytesIO()
            output_format.powerpoint(summary)
            ppt_io.seek(0)
            return StreamingResponse(ppt_io, media_type='application/vnd.openxmlformats-officedocument.presentationml.presentation', headers={"Content-Disposition": "attachment; filename=summary.pptx"})
        
    else:
            return {"error": "Unsupported format"}




